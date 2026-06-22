"""Assemble rendered slide PNGs into a 16:9 PPTX deck.

Usage:
    python build_pptx.py            # builds szu_interview.pptx
    python build_pptx.py custom.pptx # custom output name

Requires: pip install python-pptx pillow
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu


SLIDE_W = 9144000  # 1920 px at 96 dpi → EMU
SLIDE_H = 5143500  # 1080 px at 96 dpi → EMU


def build(output: Path) -> Path:
    base = Path(__file__).parent
    renders = base / "renders"

    pngs = sorted(p for p in renders.glob("*.png") if p.is_file())
    if not pngs:
        raise SystemExit(
            f"No PNGs found in {renders}. Run render.ps1 first."
        )

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # blank

    for png in pngs:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(png),
            left=Emu(0),
            top=Emu(0),
            width=Emu(SLIDE_W),
            height=Emu(SLIDE_H),
        )

    prs.save(str(output))
    return output


def main() -> None:
    # Force UTF-8 stdout so Chinese paths don't crash cp1252 consoles
    if sys.stdout.encoding and sys.stdout.encoding.lower() != "utf-8":
        try:
            sys.stdout.reconfigure(encoding="utf-8")
        except Exception:
            pass

    out_name = sys.argv[1] if len(sys.argv) > 1 else "szu_interview.pptx"
    out_path = Path(__file__).parent / out_name
    built = build(out_path)
    print(f"Built: {built.name}")
    print(f"Size: {built.stat().st_size / 1024:.1f} KB")


if __name__ == "__main__":
    main()
